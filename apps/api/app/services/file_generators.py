from pathlib import Path
from textwrap import shorten
from importlib import resources

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


BLUE = RGBColor(30, 64, 175)
DARK = RGBColor(17, 24, 39)
GRAY = RGBColor(107, 114, 128)
LIGHT_BLUE = RGBColor(239, 246, 255)


def generate_docx(path: Path, title: str, summary: str, sections: list[dict[str, str]], include_manual_checklist: bool = True) -> None:
    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(11)

    document.add_heading(title, level=1)
    document.add_paragraph(summary)
    document.add_paragraph("Конкурс: ПФКИ. Документ подготовлен как редактируемая рабочая версия.")

    for section in sections:
        document.add_heading(section["title"], level=2)
        for line in section["body"].split("\n"):
            if line.strip():
                document.add_paragraph(line.strip())

    if include_manual_checklist:
        document.add_heading("Проверить вручную перед подачей", level=2)
        for item in [
            "актуальность источников и дат",
            "точные суммы, должности, номера мероприятий и реквизиты",
            "соответствие требованиям выбранного конкурса",
        ]:
            document.add_paragraph(item, style="List Bullet")

    document.save(path)


def _register_pdf_font() -> str:
    try:
        roboto_path = resources.files("font_roboto").joinpath("files/Roboto-Regular.ttf")
        if roboto_path.is_file():
            pdfmetrics.registerFont(TTFont("LarySans", str(roboto_path)))
            return "LarySans"
    except Exception:
        pass

    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            pdfmetrics.registerFont(TTFont("LarySans", candidate))
            return "LarySans"
    return "Helvetica"


def generate_pdf(path: Path, title: str, summary: str, sections: list[dict[str, str]]) -> None:
    font = _register_pdf_font()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="LaryTitle", fontName=font, fontSize=20, leading=26, textColor=colors.HexColor("#111827"), spaceAfter=10))
    styles.add(ParagraphStyle(name="LaryHeading", fontName=font, fontSize=14, leading=18, textColor=colors.HexColor("#1E40AF"), spaceBefore=12, spaceAfter=6))
    styles.add(ParagraphStyle(name="LaryBody", fontName=font, fontSize=10.5, leading=15, textColor=colors.HexColor("#111827")))

    story = [
        Paragraph(title, styles["LaryTitle"]),
        Paragraph(summary, styles["LaryBody"]),
        Spacer(1, 6 * mm),
        Paragraph("Конкурс: ПФКИ. Рабочая версия для проверки и редактирования.", styles["LaryBody"]),
    ]
    for section in sections:
        story.append(Paragraph(section["title"], styles["LaryHeading"]))
        for line in section["body"].split("\n"):
            if line.strip():
                story.append(Paragraph(line.strip(), styles["LaryBody"]))

    doc = SimpleDocTemplate(str(path), pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm, topMargin=16 * mm, bottomMargin=16 * mm)
    doc.build(story)


def generate_pptx(path: Path, title: str, summary: str, sections: list[dict[str, str]], variant: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_plan = _presentation_plan(title, summary, sections, variant)
    for index, slide_data in enumerate(slide_plan, start=1):
        _add_slide(prs, index, slide_data["title"], slide_data["subtitle"], slide_data["bullets"], slide_data["visual"])

    prs.save(path)


def _presentation_plan(title: str, summary: str, sections: list[dict[str, str]], variant: str) -> list[dict[str, object]]:
    by_title = {item["title"]: item["body"] for item in sections}
    common = [
        {
            "title": title,
            "subtitle": "Презентация для заявки ПФКИ" if variant != "calendar_plan" else "Презентация сценарного плана",
            "bullets": [shorten(summary, width=150, placeholder="...")],
            "visual": "Ключевая идея проекта",
        },
        {
            "title": "Проблема и актуальность",
            "subtitle": "Почему проект нужен сейчас",
            "bullets": _lines(by_title.get("Актуальность", summary), 3),
            "visual": "Иллюстрация проблемы",
        },
        {
            "title": "Целевая аудитория",
            "subtitle": "Для кого создается проект",
            "bullets": _lines(by_title.get("Целевая группа", "Уточните возраст, статус, территорию и ожидаемый эффект."), 3),
            "visual": "Портрет участников",
        },
        {
            "title": "Идея решения",
            "subtitle": "Что делает проект",
            "bullets": _lines(by_title.get("Решение", "Опишите формат, механику и главный результат проекта."), 3),
            "visual": "Механика проекта",
        },
    ]

    if variant == "calendar_plan":
        middle = [
            {
                "title": "Календарный план",
                "subtitle": "Основные этапы",
                "bullets": ["Подготовка и набор участников", "Основные события и показы", "Финальная отчетность и распространение результатов"],
                "visual": "Лента времени",
            },
            {
                "title": "Сценарная логика",
                "subtitle": "Как участник проходит проект",
                "bullets": ["Вход в проект", "Серия активностей", "Финальный продукт или событие"],
                "visual": "Маршрут участника",
            },
        ]
    else:
        middle = [
            {
                "title": "Значимость для целевой группы",
                "subtitle": "Что изменится для аудитории и территории",
                "bullets": _lines(by_title.get("Значимость для целевой группы", "Покажите изменения для целевой группы и региона."), 3),
                "visual": "Ожидаемые изменения",
            },
            {
                "title": "Партнеры и поддержка",
                "subtitle": "Кто помогает проекту",
                "bullets": ["Информационная и организационная поддержка", "Экспертный вклад", "Софинансирование или имущественный вклад"],
                "visual": "Карта партнеров",
            },
        ]

    return common + middle + [
        {
            "title": "Команда и роли",
            "subtitle": "Кто отвечает за результат",
            "bullets": ["Руководитель проекта", "Координатор или организатор", "Партнеры, эксперты и исполнители"],
            "visual": "Командная схема",
        },
        {
            "title": "Результаты и показатели",
            "subtitle": "Что можно измерить",
            "bullets": ["Количество участников и зрителей", "Созданные материалы или события", "Охват, публикации, обратная связь"],
            "visual": "Показатели",
        },
        {
            "title": "Бюджетная логика",
            "subtitle": "На что направлены средства",
            "bullets": ["Оплата труда и услуги", "Материалы, площадки, техника", "Информационное сопровождение"],
            "visual": "Структура бюджета",
        },
        {
            "title": "Что проверить перед подачей",
            "subtitle": "Ручная проверка обязательна",
            "bullets": ["Источники и даты", "ФИО, суммы, реквизиты", "Календарь, бюджет и показатели"],
            "visual": "Контрольный список",
        },
    ]


def _lines(text: str, limit: int) -> list[str]:
    parts = [part.strip(" -•") for part in text.replace("\r", "\n").split("\n") if part.strip()]
    if not parts:
        parts = [text]
    return [shorten(part, width=105, placeholder="...") for part in parts[:limit]]


def _add_slide(prs: Presentation, index: int, title: str, subtitle: str, bullets: list[str], visual: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)

    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BLUE
    top_bar.line.fill.background()

    marker = slide.shapes.add_shape(1, Inches(0.55), Inches(0.55), Inches(0.7), Inches(0.7))
    marker.fill.solid()
    marker.fill.fore_color.rgb = BLUE
    marker.line.fill.background()
    marker.text_frame.text = str(index)
    marker.text_frame.paragraphs[0].font.size = PptPt(18)
    marker.text_frame.paragraphs[0].font.bold = True
    marker.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    marker.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(1.45), Inches(0.42), Inches(6.9), Inches(1.28))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(26)
    p.font.bold = True
    p.font.color.rgb = DARK

    subtitle_box = slide.shapes.add_textbox(Inches(1.45), Inches(1.7), Inches(6.8), Inches(0.45))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = PptPt(16)
    subtitle_tf.paragraphs[0].font.color.rgb = GRAY

    visual_box = slide.shapes.add_shape(1, Inches(8.7), Inches(0.85), Inches(3.85), Inches(5.4))
    visual_box.fill.solid()
    visual_box.fill.fore_color.rgb = LIGHT_BLUE
    visual_box.line.color.rgb = RGBColor(191, 219, 254)
    visual_box.text_frame.text = visual
    visual_box.text_frame.paragraphs[0].font.size = PptPt(24)
    visual_box.text_frame.paragraphs[0].font.bold = True
    visual_box.text_frame.paragraphs[0].font.color.rgb = BLUE
    visual_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    bullet_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(6.95), Inches(3.55))
    tf = bullet_box.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.font.size = PptPt(21)
        paragraph.font.color.rgb = DARK
        paragraph.space_after = PptPt(12)
        paragraph.level = 0

    footer = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.4), Inches(0.35))
    footer.text_frame.text = "Лари · рабочая презентация для редактирования"
    footer.text_frame.paragraphs[0].font.size = PptPt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = GRAY
